#!/usr/bin/env python3
"""Export existing figures to editable vector formats (EMF + PPTX).

Uses SVGs as source, converts to EMF via Inkscape, and builds PPTX slides
with EMF so you can ungroup/edit in PowerPoint.
"""

from __future__ import annotations

from pathlib import Path
import re
import shutil
import subprocess
import xml.etree.ElementTree as ET

from pptx import Presentation
from PIL import Image

Image.MAX_IMAGE_PIXELS = None


ROOT = Path("/home/ramanlab/Documents/cole")
FIG_DIR = ROOT / "Data/imagesForFigures/imagesForOptoHEXTesting"
SVG_GLOBS = [
    "*pubfig.svg",
    "*traces_bars.svg",
    "*percents.svg",
]
EMF_EXPORT_WIDTH = 2000


def _parse_length(value: str | None) -> float | None:
    if not value:
        return None
    match = re.match(r"([0-9.]+)\s*([a-z%]*)", value.strip(), re.I)
    if not match:
        return None
    num = float(match.group(1))
    unit = match.group(2).lower()
    if unit in ("", "px"):
        return num  # px at 96dpi
    if unit == "pt":
        return num * (96.0 / 72.0)
    if unit == "in":
        return num * 96.0
    if unit == "cm":
        return num * (96.0 / 2.54)
    if unit == "mm":
        return num * (96.0 / 25.4)
    return num


def _svg_aspect(svg_path: Path) -> float:
    try:
        root = ET.parse(svg_path).getroot()
    except Exception:
        return 1.0
    w = _parse_length(root.get("width"))
    h = _parse_length(root.get("height"))
    if w and h and h > 0:
        return w / h
    view_box = root.get("viewBox")
    if view_box:
        parts = [float(x) for x in view_box.strip().split()]
        if len(parts) == 4 and parts[3] != 0:
            return parts[2] / parts[3]
    return 1.0


def _export_emf(svg_path: Path, emf_path: Path, inkscape: str) -> None:
    subprocess.run(
        [
            inkscape,
            str(svg_path),
            "--export-type=emf",
            "--export-filename",
            str(emf_path),
            "--export-width",
            str(EMF_EXPORT_WIDTH),
        ],
        check=True,
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
    )


def _make_pptx(emf_path: Path, svg_path: Path, pptx_path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    aspect = _svg_aspect(svg_path)
    max_w = slide_w * 0.95
    max_h = slide_h * 0.95
    if aspect >= 1.0:
        img_w = max_w
        img_h = max_w / aspect
        if img_h > max_h:
            img_h = max_h
            img_w = max_h * aspect
    else:
        img_h = max_h
        img_w = max_h * aspect
        if img_w > max_w:
            img_w = max_w
            img_h = max_w / aspect

    left = (slide_w - img_w) / 2
    top = (slide_h - img_h) / 2
    try:
        slide.shapes.add_picture(str(emf_path), left, top, width=img_w, height=img_h)
        prs.save(str(pptx_path))
    except Exception as exc:
        print(f"[WARN] PPTX add_picture failed for {emf_path.name}: {exc}")


def main() -> None:
    inkscape = shutil.which("inkscape")
    if not inkscape:
        raise RuntimeError("Inkscape is required to export EMF.")

    svg_files: list[Path] = []
    for pattern in SVG_GLOBS:
        svg_files.extend(FIG_DIR.glob(pattern))

    if not svg_files:
        raise RuntimeError(f"No SVGs found in {FIG_DIR}")

    for svg_path in sorted(set(svg_files)):
        emf_path = svg_path.with_suffix(".emf")
        pptx_path = svg_path.with_suffix(".pptx")
        _export_emf(svg_path, emf_path, inkscape)
        _make_pptx(emf_path, svg_path, pptx_path)
        print(f"[OK] {svg_path.name} → {emf_path.name}, {pptx_path.name}")


if __name__ == "__main__":
    main()
