#!/usr/bin/env python3
"""Generate control-vs-trained PER% bar plots from binary reaction CSVs."""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager
import numpy as np
import pandas as pd
from PIL import Image

try:
    from pptx import Presentation
except Exception:  # pragma: no cover
    Presentation = None

ROOT = Path("/home/ramanlab/Documents/cole")
OUT_DIR = ROOT / "Data/imagesForFigures/imagesForOptoHEXTesting"
MATRIX_DIR = ROOT / "Results/Opto-Fly-Figures/Matrix-PER-Reactions-Model"

BAR_ORDER = [
    ("3-Octonol", "3-Octanol"),
    ("Apple Cider Vinegar", "ACV"),
    ("Benzaldehyde", "Benzaldehyde"),
    ("Citral", "Citral"),
    ("Ethyl Butyrate", "Ethyl Butyrate"),
    ("Hexanol", "Hexanol"),
    ("Linalool", "Linalool"),
]


@dataclass(frozen=True)
class PairSpec:
    title: str
    trained_dataset: str
    control_dataset: str
    trained_label: str
    control_label: str
    out_name: str


PAIR_SPECS = [
    PairSpec(
        title="PER% — Hexanol Trained vs Control",
        trained_dataset="Hex-Training",
        control_dataset="Hex-Control",
        trained_label="Trained",
        control_label="Control",
        out_name="hex_trained_vs_hex_control_percents.png",
    ),
    PairSpec(
        title="PER% — Ethyl Butyrate Trained vs Control",
        trained_dataset="EB-Training",
        control_dataset="EB-Control",
        trained_label="Trained",
        control_label="Control",
        out_name="eb_trained_vs_eb_control_percents.png",
    ),
    PairSpec(
        title="PER% — Hexanol Trained-24 vs Control-24",
        trained_dataset="Hex-Training-24",
        control_dataset="Hex-Control-24",
        trained_label="Trained-24",
        control_label="Control-24",
        out_name="hex_trained_24_vs_hex_control_24_percents.png",
    ),
    PairSpec(
        title="PER% — Hexanol Trained-36 vs Control-36",
        trained_dataset="Hex-Training-36",
        control_dataset="Hex-Control-36",
        trained_label="Trained-36",
        control_label="Control-36",
        out_name="hex_trained_36_vs_hex_control_36_percents.png",
    ),
]


def _register_arial() -> None:
    for font_path in [
        "/usr/share/fonts/truetype/msttcorefonts/Arial.ttf",
        "/usr/share/fonts/truetype/msttcorefonts/Arial_Bold.ttf",
        "/usr/share/fonts/truetype/msttcorefonts/Arial_Italic.ttf",
        "/usr/share/fonts/truetype/msttcorefonts/Arial_Bold_Italic.ttf",
        "/usr/share/fonts/truetype/msttcorefonts/arialbd.ttf",
        "/usr/share/fonts/truetype/msttcorefonts/ariali.ttf",
        "/usr/share/fonts/truetype/msttcorefonts/arialbi.ttf",
    ]:
        path = Path(font_path)
        if path.exists():
            font_manager.fontManager.addfont(str(path))


def _binary_csv_path(dataset: str) -> Path | None:
    dataset_dir = MATRIX_DIR / dataset
    candidates = sorted(dataset_dir.glob(f"binary_reactions_{dataset}*.csv"))
    if not candidates:
        return None
    for candidate in candidates:
        if "unordered" in candidate.name:
            return candidate
    return candidates[0]


def _load_dataset_rates(dataset: str) -> tuple[dict[str, float], int] | tuple[None, None]:
    csv_path = _binary_csv_path(dataset)
    if csv_path is None:
        return None, None

    try:
        df = pd.read_csv(csv_path)
    except Exception as exc:
        print(f"[WARN] Failed to read {csv_path}: {exc}")
        return None, None

    if df.empty or "odor_sent" not in df.columns or "during_hit" not in df.columns:
        print(f"[WARN] Binary reactions CSV missing usable data: {csv_path}")
        return None, None

    stats = (
        df.groupby("odor_sent")["during_hit"]
        .agg(num_reactions="sum", num_trials="size")
        .reset_index()
    )
    stats["rate"] = np.where(
        stats["num_trials"] > 0,
        (stats["num_reactions"] / stats["num_trials"]) * 100.0,
        0.0,
    )
    values = {str(row["odor_sent"]): float(row["rate"]) for _, row in stats.iterrows()}
    if {"fly", "fly_number"}.issubset(df.columns):
        n_flies = int(df[["fly", "fly_number"]].drop_duplicates().shape[0])
    else:
        n_flies = 0
    return values, n_flies


def plot_comparison(
    *,
    title: str,
    trained_values: dict[str, float],
    control_values: dict[str, float],
    trained_n: int,
    control_n: int,
    out_path: Path,
    trained_label: str,
    control_label: str,
) -> None:
    labels = [label for _, label in BAR_ORDER]
    trained = [float(trained_values.get(key, 0.0)) for key, _ in BAR_ORDER]
    control = [float(control_values.get(key, 0.0)) for key, _ in BAR_ORDER]

    x = np.arange(len(labels))
    width = 0.36

    fig, ax = plt.subplots(figsize=(11, 6))
    ax.bar(
        x - width / 2,
        control,
        width,
        color="#9e9e9e",
        edgecolor="black",
        linewidth=0.8,
        label=f"{control_label} (n = {control_n})",
    )
    ax.bar(
        x + width / 2,
        trained,
        width,
        color="#1f77b4",
        edgecolor="black",
        linewidth=0.8,
        label=f"{trained_label} (n = {trained_n})",
    )

    ax.set_ylim(0, 100)
    ax.set_ylabel("PER%", fontsize=12)
    ax.set_title(title, fontsize=14, weight="bold", loc="center", pad=6)
    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=12)
    ax.tick_params(axis="y", labelsize=14)
    ax.legend(loc="upper right", frameon=True, fontsize=14)

    _save_all(fig, out_path)
    plt.close(fig)
    print(f"[OK] Wrote {out_path}")


def _save_all(fig: plt.Figure, out_path: Path) -> None:
    out_path = Path(out_path)
    fig.savefig(out_path, bbox_inches="tight")
    svg_path = out_path.with_suffix(".svg")
    eps_path = out_path.with_suffix(".eps")
    pdf_path = out_path.with_suffix(".pdf")
    fig.savefig(svg_path, bbox_inches="tight")
    with matplotlib.rc_context(
        {
            "font.family": "Arial",
            "font.sans-serif": ["Arial"],
            "font.size": 12,
            "ps.fonttype": 42,
            "ps.useafm": False,
        }
    ):
        fig.savefig(eps_path, bbox_inches="tight")
    fig.savefig(pdf_path, bbox_inches="tight")
    _save_pptx(out_path, svg_path, out_path.with_suffix(".pptx"))


def _save_pptx(png_path: Path, svg_path: Path, pptx_path: Path) -> None:
    if Presentation is None:
        return
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    img = Image.open(png_path)
    w_px, h_px = img.size
    emu_per_px = 914400 / 96.0
    img_w = w_px * emu_per_px
    img_h = h_px * emu_per_px
    max_w = slide_w * 0.95
    max_h = slide_h * 0.95
    scale = min(max_w / img_w, max_h / img_h, 1.0)
    img_w *= scale
    img_h *= scale
    left = (slide_w - img_w) / 2
    top = (slide_h - img_h) / 2

    try:
        slide.shapes.add_picture(str(svg_path), left, top, width=img_w, height=img_h)
    except Exception:
        slide.shapes.add_picture(str(png_path), left, top, width=img_w, height=img_h)

    prs.save(str(pptx_path))


def main() -> None:
    _register_arial()
    plt.rcParams.update(
        {
            "figure.dpi": 300,
            "savefig.dpi": 300,
            "axes.spines.top": False,
            "axes.spines.right": False,
            "axes.linewidth": 0.8,
            "xtick.direction": "out",
            "ytick.direction": "out",
            "font.family": "Arial",
            "font.size": 14,
            "svg.fonttype": "none",
        }
    )
    OUT_DIR.mkdir(parents=True, exist_ok=True)

    for spec in PAIR_SPECS:
        trained_values, trained_n = _load_dataset_rates(spec.trained_dataset)
        control_values, control_n = _load_dataset_rates(spec.control_dataset)
        if trained_values is None or control_values is None:
            print(
                f"[WARN] Skipping {spec.trained_dataset} vs {spec.control_dataset}: "
                "missing binary reaction CSVs."
            )
            continue
        plot_comparison(
            title=spec.title,
            trained_values=trained_values,
            control_values=control_values,
            trained_n=trained_n or 0,
            control_n=control_n or 0,
            out_path=OUT_DIR / spec.out_name,
            trained_label=spec.trained_label,
            control_label=spec.control_label,
        )


if __name__ == "__main__":
    main()
